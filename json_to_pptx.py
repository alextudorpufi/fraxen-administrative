import json
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE

# --- CONFIGURATION ---
TEMPLATE_FILE = 'template.pptx'
JSON_FILE = 'json_output.json'
OUTPUT_FILE = 'Dora_S_Profile_Final.pptx'

# --- COLORS ---
COLOR_PURPLE = RGBColor(137, 87, 230)
COLOR_BLUE = RGBColor(0, 112, 192)
COLOR_WHITE = RGBColor(255, 255, 255)

def get_shape_by_name(slide, name):
    """Recursively finds a shape by name, even inside groups."""
    return find_in_shapes(slide.shapes, name)

def find_in_shapes(shapes, name):
    for shape in shapes:
        if shape.name == name:
            return shape
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            found = find_in_shapes(shape.shapes, name)
            if found:
                return found
    return None

def update_shape_color(shape, rgb_color):
    """Colors a shape (or all children of a group)."""
    if not shape: return
    
    # If Group, recurse
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for child in shape.shapes:
            update_shape_color(child, rgb_color)
        return

    # If simple shape, color it
    try:
        # Avoid coloring text boxes or lines
        if shape.shape_type not in [MSO_SHAPE_TYPE.LINE, MSO_SHAPE_TYPE.TEXT_BOX]:
            shape.fill.solid()
            shape.fill.fore_color.rgb = rgb_color
    except AttributeError:
        pass

def update_simple_text(shape, text, font_size=Pt(12), bold=False):
    """Updates a text box with a single style."""
    if not shape: return
    if text is None: text = ""

    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    
    run.font.color.rgb = COLOR_WHITE
    if font_size: run.font.size = font_size
    run.font.bold = bold

def update_labeled_text(shape, label, value, font_size=Pt(14)):
    """Creates the 'Label: Value' style for the sidebar."""
    if not shape: return
    
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    
    # Run 1: The Label (BOLD)
    r_label = p.add_run()
    r_label.text = label + " "
    r_label.font.bold = True
    r_label.font.color.rgb = COLOR_WHITE
    r_label.font.size = font_size
    
    # Run 2: The Value (REGULAR)
    r_value = p.add_run()
    r_value.text = value
    r_value.font.bold = False
    r_value.font.color.rgb = COLOR_WHITE
    r_value.font.size = font_size

def update_role_block(shape, role_data):
    """Rebuilds the Role block with mixed font sizes."""
    if not shape: return
    
    tf = shape.text_frame
    tf.clear()
    
    # 1. Header Line
    p = tf.paragraphs[0]
    
    # Job Title (Bold, Size 18)
    job_title = role_data.get('job_title', 'N/A')
    r_title = p.add_run()
    r_title.text = job_title + " – "
    r_title.font.bold = True
    r_title.font.color.rgb = COLOR_WHITE
    r_title.font.size = Pt(18) 
    
    # Description (Regular, Size 12)
    r_desc = p.add_run()
    r_desc.text = role_data.get('description', '')
    r_desc.font.bold = False
    r_desc.font.color.rgb = COLOR_WHITE
    r_desc.font.size = Pt(12)

    # 2. Bullets (Size 14)
    for achievement in role_data.get('achievements', []):
        p_bull = tf.add_paragraph()
        p_bull.level = 0
        
        # Arrow
        r_arrow = p_bull.add_run()
        r_arrow.text = "-> " 
        r_arrow.font.color.rgb = COLOR_WHITE
        r_arrow.font.size = Pt(14)
        
        # Text
        r_text = p_bull.add_run()
        r_text.text = achievement
        r_text.font.color.rgb = COLOR_WHITE
        r_text.font.size = Pt(14)

def update_footer_block(shape, strengths_list):
    """
    Rebuilds the Footer:
    'Core Strengths:' (Bold)
      -> List item 1 (Regular)
      -> List item 2 (Regular)
    """
    if not shape: return
    
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]

    # Run 1: The Label "Core Strengths:"
    r_label = p.add_run()
    r_label.text = "Core Strengths:"
    r_label.font.bold = True
    r_label.font.color.rgb = COLOR_WHITE
    r_label.font.size = Pt(14)

    # If there are strengths, add them on new lines
    if strengths_list:
        # Split strengths into two visual lines/groups
        midpoint = (len(strengths_list) + 1) // 2
        line1 = ", ".join(strengths_list[:midpoint])
        line2 = ", ".join(strengths_list[midpoint:])
        
        # First line of strengths
        p_line1 = tf.add_paragraph()
        p_line1.level = 0
        r_arr1 = p_line1.add_run()
        r_arr1.text = "-> "
        r_arr1.font.color.rgb = COLOR_WHITE
        r_arr1.font.size = Pt(14)
        
        r_txt1 = p_line1.add_run()
        r_txt1.text = line1
        r_txt1.font.color.rgb = COLOR_WHITE
        r_txt1.font.size = Pt(14)

        # Second line of strengths
        if line2:
            p_line2 = tf.add_paragraph()
            p_line2.level = 0
            r_arr2 = p_line2.add_run()
            r_arr2.text = "-> "
            r_arr2.font.color.rgb = COLOR_WHITE
            r_arr2.font.size = Pt(14)
            
            r_txt2 = p_line2.add_run()
            r_txt2.text = line2
            r_txt2.font.color.rgb = COLOR_WHITE
            r_txt2.font.size = Pt(14)


def main():
    print(f"Loading {TEMPLATE_FILE}...")
    try:
        with open(JSON_FILE, 'r') as f:
            data = json.load(f)
        prs = Presentation(TEMPLATE_FILE)
        slide = prs.slides[0]
    except Exception as e:
        print(f"Error loading files: {e}")
        return

    # --- 1. GENDER LOGIC (Colors) ---
    gender_val = data.get('gender', 'Female')
    
    if gender_val.strip().lower() == 'male':
        target_color = COLOR_BLUE
        print("Gender is Male -> Using BLUE.")
    else:
        target_color = COLOR_PURPLE
        print("Gender is Female -> Using PURPLE.")

    box_names = ['gender_box_main', 'gender_box_1', 'gender_box_2', 'gender_box_3', 'gender_box_4']
    for name in box_names:
        update_shape_color(get_shape_by_name(slide, name), target_color)

    # --- 2. TEXT UPDATES (Sidebar - Font 14) ---
    print("Updating Sidebar...")
    update_labeled_text(get_shape_by_name(slide, 'sidebar_gender_text'), "Gender:", gender_val, font_size=Pt(14))
    update_labeled_text(get_shape_by_name(slide, 'sidebar_sectors'), "Sector Focus:", data.get('sector_focus'), font_size=Pt(14))
    update_labeled_text(get_shape_by_name(slide, 'sidebar_location'), "Location:", data.get('location'), font_size=Pt(14))
    
    exp_text = data.get('experience_summary', '')
    update_labeled_text(get_shape_by_name(slide, 'sidebar_summary'), "Experience:", exp_text, font_size=Pt(14))

    # --- 3. TEXT UPDATES (Main Header - Font 24) ---
    print("Updating Header...")
    title_shape = get_shape_by_name(slide, 'header_title')
    update_simple_text(title_shape, data.get('title'), font_size=Pt(24), bold=True)

    # --- 4. TEXT UPDATES (Roles - Mixed Fonts) ---
    print("Updating Roles...")
    experiences = data.get('experience', [])
    for i, role_data in enumerate(experiences):
        update_role_block(get_shape_by_name(slide, f'role_{i+1}'), role_data)

    # --- 5. TEXT UPDATES (Footer - Font 14 with Label) ---
    print("Updating Footer...")
    strengths = data.get('core_strengths', [])
    update_footer_block(get_shape_by_name(slide, 'footer_strengths'), strengths)

    prs.save(OUTPUT_FILE)
    print(f"Success! Saved to {OUTPUT_FILE}")

if __name__ == "__main__":
    main()